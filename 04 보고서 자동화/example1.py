# 파이썬으로 파워포인트 편집 (python-pptx)

from pptx import Presentation
from pptx.util import Inches    # for image, chart

# 파워포인트 객체 선언
prs = Presentation()

for i in range(0, 11):
    # 슬라이드 종류 선택, 총 11개
    title_slide_layout = prs.slide_layouts[i]

    # 슬라이드 추가
    slide = prs.slides.add_slide(title_slide_layout)

prs.save('add all slides.pptx')